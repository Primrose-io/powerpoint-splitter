import os
import re
import shutil
import uuid
import zipfile
from flask import Flask, request, render_template, send_from_directory, redirect, url_for, flash, session
from werkzeug.utils import secure_filename
from pptx import Presentation

app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'dev-secret-key')
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'output'
app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50 MB limit

# Ensure upload and output directories exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

# Route to serve static files
@app.route('/static/<path:filename>')
def serve_static(filename):
    return send_from_directory('static', filename)

def extract_metadata_from_notes(notes_text):
    """
    Extracts tags from notes like [presentation1,presentation2]
    """
    if not notes_text:
        return []
    match = re.search(r'\[([^\[\]]+)\]\s*$', notes_text.strip())
    if match:
        return [tag.strip() for tag in match.group(1).split(',')]
    return []

def get_slide_tags(prs):
    """
    Returns a list where each index corresponds to a slide's tags
    """
    slide_tags = []
    for slide in prs.slides:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide and notes_slide.notes_text_frame else ""
        tags = extract_metadata_from_notes(notes_text)
        slide_tags.append(tags)
    return slide_tags

def delete_unwanted_slides(prs, keep_indices):
    """
    Deletes all slides from presentation except those in keep_indices.
    Note: Deletion must be done in reverse to avoid index shifting.
    """
    slide_id_list = prs.slides._sldIdLst
    slides_to_delete = [i for i in range(len(prs.slides)) if i not in keep_indices]
    for idx in sorted(slides_to_delete, reverse=True):
        slide_id_list.remove(slide_id_list[idx])

def process_powerpoint_file(input_path, output_dir):
    """
    Process a PowerPoint file and split it based on tags.
    Returns a list of generated files and their download paths.
    """
    # Create session-specific output directory
    os.makedirs(output_dir, exist_ok=True)
    
    # Process the PowerPoint file
    main_prs = Presentation(input_path)
    slide_tags = get_slide_tags(main_prs)
    
    # Build a tag -> [slide index] mapping
    tag_to_slide_indices = {}
    for i, tags in enumerate(slide_tags):
        for tag in tags:
            tag_to_slide_indices.setdefault(tag, []).append(i)
    
    if not tag_to_slide_indices:
        return []  # No tags found
    
    # Find slides with the '*' tag that should be included in all presentations
    universal_slides = tag_to_slide_indices.get('*', [])
    
    # Generate output files
    generated_files = []
    
    # For each tag, create filtered copy
    for tag, indices in tag_to_slide_indices.items():
        # Skip processing the '*' tag as a separate presentation
        if tag == '*':
            continue
        
        output_filename = f"{tag}.pptx"
        output_path = os.path.join(output_dir, output_filename)
        
        # Create a copy of the original file
        shutil.copyfile(input_path, output_path)
        
        # Process the copy
        tag_prs = Presentation(output_path)
        # Include both tag-specific slides and universal slides
        all_indices = sorted(set(indices + universal_slides))
        delete_unwanted_slides(tag_prs, all_indices)
        tag_prs.save(output_path)
        
        # Add to generated files list
        generated_files.append({
            'tag': tag,
            'filename': output_filename,
            'path': output_path
        })
    
    return generated_files

def create_zip_archive(files, output_dir, zip_filename):
    """Create a zip archive of all generated files"""
    zip_path = os.path.join(output_dir, zip_filename)
    
    with zipfile.ZipFile(zip_path, 'w') as zipf:
        for file_info in files:
            zipf.write(file_info['path'], arcname=file_info['filename'])
    
    return zip_path

@app.route('/')
def index():
    """Render the upload form"""
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload_file():
    # Check if a file was uploaded
    if 'file' not in request.files:
        flash('No file part')
        return redirect(request.url)
    
    file = request.files['file']
    
    # If user submits an empty form
    if file.filename == '':
        flash('No selected file')
        return redirect(request.url)
    
    # Check if the file is a PowerPoint file
    if file and file.filename.endswith(('.ppt', '.pptx')):
        # Create unique session ID
        session_id = str(uuid.uuid4())
        session['session_id'] = session_id
        
        # Create session-specific directories
        session_upload_dir = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
        session_output_dir = os.path.join(app.config['OUTPUT_FOLDER'], session_id)
        
        os.makedirs(session_upload_dir, exist_ok=True)
        os.makedirs(session_output_dir, exist_ok=True)
        
        # Save the uploaded file
        filename = secure_filename(file.filename)
        file_path = os.path.join(session_upload_dir, filename)
        file.save(file_path)
        
        # Process the file
        try:
            generated_files = process_powerpoint_file(file_path, session_output_dir)
            
            if not generated_files:
                flash('No tags found in the presentation. Please ensure slides contain tags in notes like [tag1,tag2].')
                return redirect(url_for('index'))
            
            # Create a zip archive of all files
            zip_filename = 'all_presentations.zip'
            zip_path = create_zip_archive(generated_files, session_output_dir, zip_filename)
            
            session['generated_files'] = generated_files
            session['zip_filename'] = zip_filename
            
            return redirect(url_for('show_results'))
        
        except Exception as e:
            flash(f'Error processing file: {str(e)}')
            return redirect(url_for('index'))
    
    else:
        flash('Invalid file type. Please upload a PowerPoint file (.ppt or .pptx)')
        return redirect(url_for('index'))

@app.route('/results')
def show_results():
    """Show the processed results and download links"""
    if 'session_id' not in session or 'generated_files' not in session:
        flash('No processed files found. Please upload a PowerPoint file first.')
        return redirect(url_for('index'))
    
    return render_template('results.html', 
                          files=session['generated_files'],
                          zip_filename=session.get('zip_filename'))

@app.route('/download/<filename>')
def download_file(filename):
    """Download a processed file"""
    if 'session_id' not in session:
        flash('Session expired. Please upload your file again.')
        return redirect(url_for('index'))
    
    session_output_dir = os.path.join(app.config['OUTPUT_FOLDER'], session['session_id'])
    return send_from_directory(session_output_dir, filename, as_attachment=True)

@app.route('/cleanup', methods=['POST'])
def cleanup():
    """Clean up session files"""
    if 'session_id' in session:
        session_id = session['session_id']
        session_upload_dir = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
        session_output_dir = os.path.join(app.config['OUTPUT_FOLDER'], session_id)
        
        # Remove directories if they exist
        if os.path.exists(session_upload_dir):
            shutil.rmtree(session_upload_dir)
        
        if os.path.exists(session_output_dir):
            shutil.rmtree(session_output_dir)
        
        # Clear session data
        session.clear()
    
    return redirect(url_for('index'))

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5001)), debug=False)