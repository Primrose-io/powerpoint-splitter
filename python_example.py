import os
import re
import shutil
from pptx import Presentation

# Static path to your main presentation and output folder
MAIN_PPTX = "/Users/simon/repos/powerpoint_splitter/sample/Introduction to Primrose.pptx"
OUTPUT_DIR = "/Users/simon/repos/powerpoint_splitter/sample/output"

def extract_metadata_from_notes(notes_text):
    """
    Extracts tags from notes like [presentation1,presentation2]
    """
    if not notes_text:
        return []
    match = re.search(r'\[([^\[\]]+)\]\s*$', notes_text.strip())
    if match:
        return [tag.strip() for tag in match.group(1).split(',')]
    return []

def get_slide_tags(prs):
    """
    Returns a list where each index corresponds to a slide's tags
    """
    slide_tags = []
    for slide in prs.slides:
        notes_slide = slide.notes_slide
        notes_text = notes_slide.notes_text_frame.text if notes_slide and notes_slide.notes_text_frame else ""
        tags = extract_metadata_from_notes(notes_text)
        slide_tags.append(tags)
    return slide_tags

def delete_unwanted_slides(prs, keep_indices):
    """
    Deletes all slides from presentation except those in keep_indices.
    Note: Deletion must be done in reverse to avoid index shifting.
    """
    slide_id_list = prs.slides._sldIdLst
    slides_to_delete = [i for i in range(len(prs.slides)) if i not in keep_indices]
    for idx in sorted(slides_to_delete, reverse=True):
        slide_id_list.remove(slide_id_list[idx])

def create_mini_decks_by_deletion():
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    # First pass: get tags from the main deck
    main_prs = Presentation(MAIN_PPTX)
    slide_tags = get_slide_tags(main_prs)

    # Build a tag -> [slide index] mapping
    tag_to_slide_indices = {}
    for i, tags in enumerate(slide_tags):
        for tag in tags:
            tag_to_slide_indices.setdefault(tag, []).append(i)
    
    # Find slides with the '*' tag that should be included in all presentations
    universal_slides = tag_to_slide_indices.get('*', [])

    # For each tag, create filtered copy
    for tag, indices in tag_to_slide_indices.items():
        # Skip processing the '*' tag as a separate presentation if desired
        if tag == '*':
            continue
            
        copy_path = os.path.join(OUTPUT_DIR, f"{tag}.pptx")
        shutil.copyfile(MAIN_PPTX, copy_path)

        tag_prs = Presentation(copy_path)
        # Include both tag-specific slides and universal slides marked with '*'
        all_indices = sorted(set(indices + universal_slides))
        delete_unwanted_slides(tag_prs, all_indices)
        tag_prs.save(copy_path)
        print(f"✔ Created: {copy_path}")

# Run the script
if __name__ == "__main__":
    create_mini_decks_by_deletion()